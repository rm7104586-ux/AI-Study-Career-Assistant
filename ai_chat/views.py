import json
import os
import re
from pathlib import Path

from google import genai
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response
from rest_framework.views import APIView

from .models import Activity


# =========================================================
# Gemini Configuration
# =========================================================

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

gemini_client = None

if GEMINI_API_KEY:
    gemini_client = genai.Client(
        api_key=GEMINI_API_KEY
    )

GEMINI_MODEL = "gemini-3.5-flash"


# =========================================================
# Gemini Helper
# =========================================================

def generate_ai_response(prompt):
    if not gemini_client:
        raise ValueError(
            "GEMINI_API_KEY is not configured."
        )

    response = gemini_client.models.generate_content(
        model=GEMINI_MODEL,
        contents=prompt,
    )

    if not response.text:
        raise ValueError(
            "Gemini returned an empty response."
        )

    return response.text.strip()


# =========================================================
# File Text Extraction
# =========================================================

def extract_text_from_file(uploaded_file):
    filename = uploaded_file.name.lower()
    extension = Path(filename).suffix

    # PDF
    if extension == ".pdf":
        reader = PdfReader(uploaded_file)

        text = ""

        for page in reader.pages:
            page_text = page.extract_text()

            if page_text:
                text += page_text + "\n"

        return text.strip(), len(reader.pages)

    # DOCX
    elif extension == ".docx":
        document = Document(uploaded_file)

        paragraphs = []

        for paragraph in document.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(paragraph.text)

        text = "\n".join(paragraphs)

        return text.strip(), None

    # TXT
    elif extension == ".txt":
        text = uploaded_file.read().decode(
            "utf-8",
            errors="ignore"
        )

        return text.strip(), None

    # PPTX
    elif extension == ".pptx":
        presentation = Presentation(uploaded_file)

        slides_text = []

        for slide in presentation.slides:
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        slide_text.append(shape.text)

            if slide_text:
                slides_text.append(
                    "\n".join(slide_text)
                )

        text = "\n\n".join(slides_text)

        return text.strip(), len(presentation.slides)

    else:
        raise ValueError(
            "Unsupported file type. "
            "Please upload PDF, DOCX, TXT, or PPTX."
        )


# =========================================================
# AI Study Chat
# =========================================================

class ChatView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        message = request.data.get("message")

        if not message:
            return Response(
                {"error": "Message is required."},
                status=400
            )

        try:
            prompt = f"""
You are an AI Study and Career Assistant.

Help a college student understand the following question.

Student question:
{message}

Give a clear, accurate and useful answer.

Use simple language when appropriate.
Use headings, bullet points and examples when helpful.
"""

            reply = generate_ai_response(prompt)

            Activity.objects.create(
                user=request.user,
                activity_type="chat",
                title="Asked a question in AI Study Chat",
                description=message,
            )

            return Response({
                "reply": reply
            })

        except Exception as error:
            return Response(
                {
                    "error": (
                        "AI service error: "
                        f"{str(error)}"
                    )
                },
                status=500
            )


# =========================================================
# Notes Analyzer
# =========================================================

class NotesUploadView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        uploaded_file = request.FILES.get("file")

        if not uploaded_file:
            return Response(
                {"error": "No study file was uploaded."},
                status=400
            )

        allowed_extensions = [
            ".pdf",
            ".docx",
            ".txt",
            ".pptx",
        ]

        extension = Path(
            uploaded_file.name.lower()
        ).suffix

        if extension not in allowed_extensions:
            return Response(
                {
                    "error": (
                        "Unsupported file type. "
                        "Allowed: PDF, DOCX, TXT, PPTX."
                    )
                },
                status=400
            )

        try:
            text, page_count = extract_text_from_file(
                uploaded_file
            )

            if not text:
                return Response(
                    {
                        "error": (
                            "No readable text was found "
                            "in this file."
                        )
                    },
                    status=400
                )

            ai_text = text[:30000]

            prompt = f"""
You are an AI Study Assistant.

Analyze this study material:

-------------------------
{ai_text}
-------------------------

Create useful study notes.

Use exactly these sections:

SUMMARY
Give a clear summary.

KEY POINTS
List the most important points.

IMPORTANT CONCEPTS
List important concepts, definitions, formulas or ideas.

REVISION TOPICS
List what the student should revise.

EXAM QUESTIONS
Create 5 practice questions from the material.

Use only information supported by the uploaded material.
Do not invent information.
Use clear language suitable for a college student.
"""

            summary = generate_ai_response(prompt)

            Activity.objects.create(
                user=request.user,
                activity_type="notes",
                title="Analyzed study material",
                description=uploaded_file.name,
            )

            words = text.split()

            return Response({
                "filename": uploaded_file.name,
                "file_type": extension.upper().replace(
                    ".",
                    ""
                ),
                "pages": page_count,
                "word_count": len(words),
                "text": text,
                "summary": summary,
            })

        except Exception as error:
            return Response(
                {"error": str(error)},
                status=500
            )


# =========================================================
# Quiz Generator
# =========================================================

class QuizView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        action = request.data.get(
            "action",
            "generate"
        )

        # -------------------------------------------------
        # Save completed quiz
        # -------------------------------------------------

        if action == "complete":

            topic = request.data.get("topic")
            score = request.data.get("score")
            total_questions = request.data.get(
                "totalQuestions"
            )

            if not topic:
                return Response(
                    {"error": "Quiz topic is required."},
                    status=400
                )

            if score is None or total_questions is None:
                return Response(
                    {
                        "error": (
                            "Score and totalQuestions "
                            "are required."
                        )
                    },
                    status=400
                )

            Activity.objects.create(
                user=request.user,
                activity_type="quiz",
                title="Completed a quiz",
                description=(
                    f"{topic} - "
                    f"Score: {score}/{total_questions}"
                ),
            )

            return Response({
                "message": (
                    "Quiz activity saved successfully."
                )
            })

        # -------------------------------------------------
        # Generate AI quiz
        # -------------------------------------------------

        topic = request.data.get("topic")

        question_count = request.data.get(
            "questionCount",
            5
        )

        difficulty = request.data.get(
            "difficulty",
            "Medium"
        )

        if not topic:
            return Response(
                {"error": "Topic is required."},
                status=400
            )

        try:
            question_count = int(question_count)

        except (TypeError, ValueError):
            question_count = 5

        question_count = max(
            1,
            min(question_count, 20)
        )

        prompt = f"""
Create a multiple-choice quiz for a college student.

Topic:
{topic}

Difficulty:
{difficulty}

Number of questions:
{question_count}

Create exactly {question_count} questions.

Each question must have:
- question
- exactly four options
- one correct answer

Return ONLY valid JSON.

Use this format:

{{
  "topic": "{topic}",
  "difficulty": "{difficulty}",
  "questions": [
    {{
      "question": "Question text",
      "options": [
        "Option A",
        "Option B",
        "Option C",
        "Option D"
      ],
      "answer": "Option A"
    }}
  ]
}}

Do not include markdown code fences.
"""

        try:
            ai_response = generate_ai_response(
                prompt
            )

            cleaned_response = ai_response.strip()

            if cleaned_response.startswith("```"):
                cleaned_response = re.sub(
                    r"^```(?:json)?\s*",
                    "",
                    cleaned_response
                )

                cleaned_response = re.sub(
                    r"\s*```$",
                    "",
                    cleaned_response
                )

            quiz_data = json.loads(
                cleaned_response
            )

            questions = quiz_data.get(
                "questions",
                []
            )

            if not questions:
                raise ValueError(
                    "Gemini did not return quiz questions."
                )

            validated_questions = []

            for question in questions:

                question_text = question.get(
                    "question"
                )

                options = question.get(
                    "options"
                )

                answer = question.get(
                    "answer"
                )

                if not question_text:
                    continue

                if not isinstance(options, list):
                    continue

                if len(options) != 4:
                    continue

                if answer not in options:
                    continue

                validated_questions.append({
                    "question": question_text,
                    "options": options,
                    "answer": answer,
                })

            validated_questions = (
                validated_questions[:question_count]
            )

            if not validated_questions:
                raise ValueError(
                    "No valid quiz questions were returned."
                )

            return Response({
                "topic": quiz_data.get(
                    "topic",
                    topic
                ),
                "difficulty": quiz_data.get(
                    "difficulty",
                    difficulty
                ),
                "questionCount": len(
                    validated_questions
                ),
                "questions": validated_questions,
            })

        except json.JSONDecodeError:
            return Response(
                {
                    "error": (
                        "Gemini returned an invalid "
                        "quiz format."
                    )
                },
                status=500
            )

        except Exception as error:
            return Response(
                {"error": str(error)},
                status=500
            )


# =========================================================
# Resume Analyzer
# =========================================================

class ResumeUploadView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        uploaded_file = request.FILES.get("file")

        if not uploaded_file:
            return Response(
                {"error": "No resume file was uploaded."},
                status=400
            )

        if not uploaded_file.name.lower().endswith(
            ".pdf"
        ):
            return Response(
                {"error": "Only PDF resumes are allowed."},
                status=400
            )

        try:
            reader = PdfReader(uploaded_file)

            text = ""

            for page in reader.pages:
                page_text = page.extract_text()

                if page_text:
                    text += page_text + "\n"

            text = text.strip()

            if not text:
                return Response(
                    {
                        "error": (
                            "No readable text was found "
                            "in this resume."
                        )
                    },
                    status=400
                )

            # ---------------------------------------------
            # Existing rule-based checks
            # ---------------------------------------------

            lower_text = text.lower()

            section_keywords = {
                "Contact Information": [
                    "email",
                    "phone",
                    "mobile",
                    "linkedin",
                ],
                "Education": [
                    "education",
                    "b.tech",
                    "btech",
                    "bachelor",
                    "degree",
                    "university",
                    "college",
                ],
                "Skills": [
                    "skills",
                    "technical skills",
                    "technologies",
                ],
                "Projects": [
                    "projects",
                    "project",
                ],
                "Experience": [
                    "experience",
                    "internship",
                    "intern",
                    "work experience",
                ],
                "Certifications": [
                    "certification",
                    "certifications",
                    "certificate",
                ],
                "Achievements": [
                    "achievement",
                    "achievements",
                    "awards",
                ],
            }

            detected_sections = []

            for section, keywords in (
                section_keywords.items()
            ):
                if any(
                    keyword in lower_text
                    for keyword in keywords
                ):
                    detected_sections.append(section)

            skill_list = [
                "python",
                "java",
                "javascript",
                "typescript",
                "c",
                "c++",
                "sql",
                "html",
                "css",
                "react",
                "node.js",
                "nodejs",
                "django",
                "flask",
                "postgresql",
                "mysql",
                "mongodb",
                "git",
                "github",
                "machine learning",
                "deep learning",
                "nlp",
                "tensorflow",
                "pytorch",
                "docker",
                "aws",
            ]

            detected_skills = []

            for skill in skill_list:
                if skill in lower_text:
                    detected_skills.append(skill)

            email_found = bool(
                re.search(
                    r"[A-Za-z0-9._%+-]+"
                    r"@[A-Za-z0-9.-]+\.[A-Za-z]{2,}",
                    text
                )
            )

            phone_found = bool(
                re.search(
                    r"(?:\+91[\s-]?)?[6-9]\d{9}\b",
                    text
                )
            )

            projects_found = any(
                keyword in lower_text
                for keyword in [
                    "projects",
                    "project",
                ]
            )

            education_found = any(
                keyword in lower_text
                for keyword in [
                    "education",
                    "b.tech",
                    "btech",
                    "bachelor",
                    "degree",
                    "college",
                    "university",
                ]
            )

            experience_found = any(
                keyword in lower_text
                for keyword in [
                    "experience",
                    "internship",
                    "intern",
                    "work experience",
                ]
            )

            skills_found = any(
                keyword in lower_text
                for keyword in [
                    "skills",
                    "technical skills",
                    "technologies",
                ]
            )

            required_sections = [
                "Education",
                "Skills",
                "Projects",
                "Experience",
            ]

            missing_sections = [
                section
                for section in required_sections
                if section not in detected_sections
            ]

            score = 0

            if email_found:
                score += 10

            if phone_found:
                score += 10

            if education_found:
                score += 15

            if skills_found:
                score += 20

            if projects_found:
                score += 20

            if experience_found:
                score += 15

            if "Certifications" in detected_sections:
                score += 5

            if "Achievements" in detected_sections:
                score += 5

            score = min(score, 100)

            suggestions = []

            if not email_found:
                suggestions.append(
                    "Add a professional email address."
                )

            if not phone_found:
                suggestions.append(
                    "Add a valid phone number."
                )

            if not education_found:
                suggestions.append(
                    "Add an Education section with "
                    "your degree and institution."
                )

            if not skills_found:
                suggestions.append(
                    "Add a clear Skills section with "
                    "relevant technical skills."
                )

            if not projects_found:
                suggestions.append(
                    "Add projects with technologies "
                    "and measurable outcomes."
                )

            if not experience_found:
                suggestions.append(
                    "Add internships, work experience, "
                    "or relevant practical experience."
                )

            if "Certifications" not in detected_sections:
                suggestions.append(
                    "Consider adding relevant certifications."
                )

            if "Achievements" not in detected_sections:
                suggestions.append(
                    "Consider adding important academic "
                    "or professional achievements."
                )

            if not suggestions:
                suggestions.append(
                    "Your resume contains the major "
                    "sections detected by this analyzer."
                )

            # ---------------------------------------------
            # Real Gemini resume analysis
            # ---------------------------------------------

            ai_text = text[:30000]

            prompt = f"""
You are an AI Resume and Career Assistant.

Analyze the following college student's resume.

-------------------------
{ai_text}
-------------------------

Provide a professional but student-friendly analysis.

Use exactly these sections:

OVERALL FEEDBACK
Give a concise assessment of the resume.

STRENGTHS
List the strongest parts of the resume.

WEAKNESSES
List areas that need improvement.

SKILLS FEEDBACK
Comment on the technical and soft skills shown.

PROJECT FEEDBACK
Explain how the projects could be presented more effectively.

ATS RECOMMENDATIONS
Give practical suggestions to improve ATS readability.

CAREER RECOMMENDATIONS
Suggest realistic improvements for internships and entry-level roles.

ACTION PLAN
Give the top 5 changes the student should make first.

Do not invent experience, skills, projects, education,
or achievements that are not present in the resume.
Base the analysis only on the uploaded resume.
"""

            ai_analysis = generate_ai_response(prompt)

            # ---------------------------------------------
            # Combined analysis
            # ---------------------------------------------

            word_count = len(text.split())

            basic_analysis = (
                "Basic Resume Analysis\n\n"
                f"Resume Score: {score}/100\n\n"
                f"Detected Sections: "
                f"{', '.join(detected_sections) if detected_sections else 'None'}\n\n"
                f"Detected Skills: "
                f"{', '.join(detected_skills) if detected_skills else 'None'}\n\n"
                f"Missing Recommended Sections: "
                f"{', '.join(missing_sections) if missing_sections else 'None'}\n\n"
                "Basic Suggestions:\n"
                + "\n".join(
                    f"{index + 1}. {suggestion}"
                    for index, suggestion in enumerate(
                        suggestions
                    )
                )
            )

            final_analysis = (
                basic_analysis
                + "\n\n"
                + "=" * 50
                + "\n\n"
                + "GEMINI AI ANALYSIS\n\n"
                + ai_analysis
            )

            Activity.objects.create(
                user=request.user,
                activity_type="resume",
                title="Resume analyzed",
                description=(
                    f"{uploaded_file.name} - "
                    f"Score: {score}/100"
                ),
            )

            return Response({
                "filename": uploaded_file.name,
                "pages": len(reader.pages),
                "word_count": word_count,
                "score": score,
                "detected_sections": detected_sections,
                "detected_skills": detected_skills,
                "missing_sections": missing_sections,
                "suggestions": suggestions,
                "analysis": final_analysis,
                "ai_analysis": ai_analysis,
                "text": text,
            })

        except Exception as error:
            return Response(
                {"error": str(error)},
                status=500
            )


# =========================================================
# Activity
# =========================================================

class ActivityListView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        activities = Activity.objects.filter(
            user=request.user
        ).order_by("-created_at")[:10]

        data = []

        for activity in activities:
            data.append({
                "id": activity.id,
                "type": activity.activity_type,
                "title": activity.title,
                "description": activity.description,
                "time": activity.created_at.strftime(
                    "%d %b %Y, %I:%M %p"
                ),
            })

        return Response(data)